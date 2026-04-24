import logging
from datetime import datetime
from contextlib import contextmanager
import locale
from pptx import Presentation
from pptx.util import Inches


LOGGER = logging.getLogger(__name__)


@contextmanager
def locale_block(local_name: str):
    lc_var: int = locale.LC_ALL
    org_local = locale.getlocale()
    try:
        yield locale.setlocale(lc_var, local_name)
    finally:
        locale.setlocale(lc_var, org_local)


class DocumentPPT:

    def __init__(self, param: dict, post_param: dict):
        self.param: dict = param
        self.post_param: dict = post_param

        self.slides = []

        self.prs = Presentation(self.post_param['ppt_template_abs_file_path'])
        self.w = 13.333
        self.h = 7.5
        self.header = 0.7

        # use custom template (and add more placeholders for you slides)
        self.slide_type = {
                'sld_title': 0,
                'sld_layout_title_and_content': 1,  # this is a standard slide for content, used everywhere
                'sld_section_header': 2,
                'sld_side_by_side_textboxes': 3,
                'sld_comparizon': 4,
                'sld_title_only': 5,
                'sld_blank': 6,
                'sld_content_with_caption': 7,
                'sld_picture_with_caption': 8,
                'sld_comparizon_no_pictures': 11
            }

    def add_new_slide(self, slide_type: str, slide_config=None):
        try:
            if slide_config is None:
                slide_config = []

            if slide_type == 'title':
                self.add_title_slide()
            elif slide_type == 'section_header':
                self.add_slide_section_header(slide_config)
            elif slide_type == 'bullets':
                self.add_bullets_slide(slide_config)
            elif slide_type == 'images':
                self.add_images_slide(slide_config)
            else:
                raise ValueError(f"Unknown slide type '{slide_type}'")

        except KeyError as _err:
            LOGGER.error(f"POS KeyError: {str(_err)}")
        except ValueError as _err:
            LOGGER.error(f"POS ValueError: {str(_err)}")
        except Exception as _err:
            LOGGER.error(f"POS Exception: {str(_err)}")
        else:
            return
        raise Exception("Error in add_new_slide")

    @staticmethod
    def test_and_print_placeholders(placeholders):
        for i, placeholder in enumerate(placeholders):
            print(f'i={i}, idx={placeholder.placeholder_format.idx}, name={placeholder.name}, '
                  f'type={placeholder.placeholder_format.type}')

    def add_title_slide(self):
        try:
            layout = self.prs.slide_layouts[self.slide_type['sld_title']]
            slide = self.prs.slides.add_slide(layout)
            title = slide.shapes.title
            # a placeholder is a srv_pre-formatted container into which content can be placed.
            subtitle = slide.placeholders[1]
            # item_title = slide.placeholders[0]
            with locale_block('zh'):
                date = datetime.today().strftime(u"%Y年 %m月 %d日")
            title.text = self.post_param['ppt_report_number']
            subtitle.text = self.post_param['ppt_report_name']

            # You can add a text placeholder (like in GUI) by setting text and coordinates
            left = top = width = height = Inches(1)
            tx_box = slide.shapes.add_textbox(left, top, width, height)
            tf = tx_box.text_frame
            tf.text = self.post_param['ppt_title']
            self.prs.save(self.post_param['local_ppt_abs_file_path'])
            self.slides.append(slide)

        except KeyError as _err:
            LOGGER.error(f"POS KeyError: {str(_err)}")
        except ValueError as _err:
            LOGGER.error(f"POS ValueError: {str(_err)}")
        except Exception as _err:
            LOGGER.error(f"POS Exception: {str(_err)}")
        else:
            return
        raise Exception("Error in add_title_slide")

    def add_slide_section_header(self, slide_config: list):
        txt = slide_config[0]  # Add text as argument
        assert isinstance(txt, str)
        layout = self.prs.slide_layouts[self.slide_type['sld_section_header']]
        slide = self.prs.slides.add_slide(layout)
        title_header = slide.shapes.title
        title_header.text = txt
        self.prs.save(self.post_param['local_ppt_abs_file_path'])
        self.slides.append(slide)
        # check amount of placeholders on the slide
        # for shape in slide.placeholders:
        #     print('%d %s' % (shape.placeholder_format.idx, shape.name))

    def add_bullets_slide(self, slide_config: list):
        assert len(slide_config) == 2, 'Number of items in slide_config must be 2'

        title, list_of_bullets = slide_config  # Title, None

        assert isinstance(title, str)
        assert isinstance(list_of_bullets, list)
        assert len(list_of_bullets) % 2 == 0, 'Number of items in list of bullets must be even'

        layout = self.prs.slide_layouts[self.slide_type['sld_layout_title_and_content']]
        slide = self.prs.slides.add_slide(layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = title

        tf = body_shape.text_frame
        for i in range(0, len(list_of_bullets), 2):
            level = list_of_bullets[i]
            string = list_of_bullets[i + 1]
            assert isinstance(level, int), 'Level of bullet must be integer number'
            assert 0 <= level <= 4, 'Level of bullet must be between 0 and 4'
            assert isinstance(string, str)
            if i == 0:
                tf.text = string
            else:
                p = tf.add_paragraph()
                p.text = string
                p.level = level

        self.prs.save(self.post_param['local_ppt_abs_file_path'])
        self.slides.append(slide)

    def add_images_slide(self, args) -> None:
        try:
            title, cols, rows = args
            title: str
            cols: list[list[str]]
            rows: list[list[str]]
            layout = self.prs.slide_layouts[self.slide_type['sld_comparizon_no_pictures']]
            slide = self.prs.slides.add_slide(layout)
            placeholders = slide.placeholders

            placeholders[0].text = title

            height_header = 0.8
            height_footer = 0.556

            height_canvas = self.h - height_header - height_footer

            left = 0

            for i, col in enumerate(cols):
                height = height_canvas / len(col)
                top = height_header

                for img_path in col:
                    slide.shapes.add_picture(img_path, Inches(left), Inches(top),
                                             width=Inches(height), height=Inches(height))
                    top += height

                left += height

            factors = [1 / len(row) for row in rows]
            sum_factors = sum(factors)
            factors = [factor / sum_factors for factor in factors]  # normalize

            top = height_header
            for i, row in enumerate(rows):
                width = min(height_canvas * factors[i], (self.w - left) / len(row))

                row_left = left
                for img_path in row:
                    slide.shapes.add_picture(img_path, Inches(row_left), Inches(top),
                                             width=Inches(width), height=Inches(width))
                    row_left += width

                top += width

            self.prs.save(self.post_param['local_ppt_abs_file_path'])
            self.slides.append(slide)

        except KeyError as _err:
            LOGGER.error(f"POS KeyError: {str(_err)}")
        except ValueError as _err:
            LOGGER.error(f"POS ValueError: {str(_err)}")
        except Exception as _err:
            LOGGER.error(f"POS Exception: {str(_err)}")
        else:
            return
        raise RuntimeError("Error in add_images_slide")
